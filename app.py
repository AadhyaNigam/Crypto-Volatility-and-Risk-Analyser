from flask import Flask, jsonify, render_template, send_file, make_response
import requests
import pandas as pd
from datetime import datetime
import os
import io
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
from apscheduler.schedulers.background import BackgroundScheduler

app = Flask(__name__)

CSV_FILE = "crypto_data.csv"
TOP_N = 5
last_snapshot = []


# ---------------------------
# FETCH LIVE MARKET DATA
# ---------------------------
def fetch_crypto_data():
    global last_snapshot

    url = "https://api.coingecko.com/api/v3/coins/markets"
    params = {
        "vs_currency": "usd",
        "order": "market_cap_desc",
        "per_page": TOP_N,
        "page": 1,
        "sparkline": False
    }

    try:
        response = requests.get(url, params=params, timeout=10)
        response.raise_for_status()
        data = response.json()
    except:
        return last_snapshot  # fallback

    rows = []
    now = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

    for coin in data:
        rows.append({
            "time": now,
            "name": coin.get("name", ""),
            "symbol": coin.get("symbol", "").upper(),
            "price": coin.get("current_price", 0.0),
            "change": coin.get("price_change_percentage_24h") or 0.0,
            "volume": coin.get("total_volume", 0),
        })

    last_snapshot = rows

    df_new = pd.DataFrame(rows)
    if not os.path.exists(CSV_FILE):
        df_new.to_csv(CSV_FILE, index=False)
    else:
        df_new.to_csv(CSV_FILE, mode='a', header=False, index=False)

    return rows


# ---------------------------
# BACKGROUND LOGGING
# ---------------------------
scheduler = BackgroundScheduler()
scheduler.add_job(fetch_crypto_data, 'interval', seconds=30)
scheduler.start()


# ---------------------------
# ROUTES
# ---------------------------
@app.route("/")
def home():
    return render_template("index.html")


@app.route("/data")
def data():
    result = fetch_crypto_data()
    return jsonify(result)


# ---------------------------
# BTC VOLATILITY IMAGE CHART
# ---------------------------
@app.route("/chart/<symbol>")
def chart(symbol):
    if not os.path.exists(CSV_FILE):
        return "No data yet", 404

    df = pd.read_csv(CSV_FILE, parse_dates=["time"])
    df = df[df["symbol"].str.upper() == symbol.upper()].copy()

    if df.empty:
        return "No data for this symbol", 404

    df.sort_values("time", inplace=True)
    df["return"] = df["price"].pct_change()
    df["rolling_vol"] = df["return"].rolling(10).std() * (252**0.5)

    fig, ax = plt.subplots(figsize=(8,4))
    ax.plot(df["time"], df["rolling_vol"])
    ax.set_title(f"{symbol.upper()} Volatility")
    ax.set_ylabel("Volatility")
    ax.set_xlabel("Time")
    fig.autofmt_xdate()

    buf = io.BytesIO()
    plt.savefig(buf, format="png", bbox_inches="tight")
    plt.close(fig)
    buf.seek(0)

    return send_file(buf, mimetype="image/png")


# ---------------------------
# BTC 7-DAY PRICE HISTORY FOR CHART.JS
# ---------------------------
@app.route("/btc_history")
def btc_history():
    url = "https://api.coingecko.com/api/v3/coins/bitcoin/market_chart"
    params = {"vs_currency": "usd", "days": 7, "interval": "daily"}

    try:
        response = requests.get(url, params=params, timeout=10)
        response.raise_for_status()
        data = response.json()
    except:
        return jsonify([])

    prices = data.get("prices", [])

    cleaned = [
        {
            "date": pd.to_datetime(p[0], unit="ms").strftime("%Y-%m-%d"),
            "price": round(p[1], 2)
        }
        for p in prices
    ]

    return jsonify(cleaned)


# ---------------------------
# PPTX EXPORT
# ---------------------------
@app.route("/pptx")
def ppt():
    if not os.path.exists(CSV_FILE):
        return "No data yet", 404

    df = pd.read_csv(CSV_FILE)
    latest_time = df["time"].max()
    latest_df = df[df["time"] == latest_time][['name','symbol','price','change','volume']]

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"Crypto Report - {latest_time}"

    rows = latest_df.shape[0] + 1
    cols = latest_df.shape[1]

    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(9), Inches(1.5)).table

    for i, col in enumerate(latest_df.columns):
        table.cell(0, i).text = col

    for r, (_, row) in enumerate(latest_df.iterrows(), start=1):
        for c, val in enumerate(row):
            table.cell(r, c).text = str(val)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    response = make_response(buf.getvalue())
    response.headers.set("Content-Disposition", "attachment", filename="crypto_report.pptx")
    response.headers.set("Content-Type", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
    return response


# ---------------------------
# RUN APP
# ---------------------------
if __name__ == "__main__":
    app.run(debug=True, use_reloader=False)